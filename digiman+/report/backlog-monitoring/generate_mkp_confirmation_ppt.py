# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x6B, 0x8F, 0x6E)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xF9, 0xA8, 0x25)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = GREEN
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = GRAY
    # underline bar
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(12.1), Pt(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

def add_item(slide, top, item_text, why_text, accent=False):
    # left bar
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(top), Pt(4), Inches(1.05) if not accent else Inches(1.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT if accent else GREEN
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.85), Inches(top - 0.03), Inches(11.9), Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = "Perlu dikonfirmasi: "
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = DARK
    r2 = p1.add_run()
    r2.text = item_text
    r2.font.size = Pt(15)
    r2.font.color.rgb = DARK

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r3 = p2.add_run()
    r3.text = "Kenapa dibutuhkan: "
    r3.font.bold = True
    r3.font.italic = True
    r3.font.size = Pt(13)
    r3.font.color.rgb = GREEN
    r4 = p2.add_run()
    r4.text = why_text
    r4.font.italic = True
    r4.font.size = Pt(13)
    r4.font.color.rgb = GRAY

def add_table(slide, left, top, width, height, headers, rows):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF3, 0xF8, 0xF3) if i % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.color.rgb = DARK
    return shape

def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY


# ===== SLIDE 1 - TITLE =====
s = add_slide()
box = slide_title = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Hal-Hal yang Perlu Dikonfirmasi ke Tim MKP"
r.font.size = Pt(40)
r.font.bold = True
r.font.color.rgb = GREEN
p2 = tf.add_paragraph()
p2.space_before = Pt(14)
r2 = p2.add_run()
r2.text = "Backlog Monitoring Dashboard — Adaptasi SP usp_iams_backlog_monitoring untuk Tenant MKP"
r2.font.size = Pt(18)
r2.font.color.rgb = GRAY
p3 = tf.add_paragraph()
p3.space_before = Pt(30)
r3 = p3.add_run()
r3.text = "Prepared by Okky Rizki R  |  10 Juni 2026"
r3.font.size = Pt(13)
r3.font.color.rgb = GRAY
bar = s.shapes.add_shape(1, Inches(1), Inches(2.45), Inches(4), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()


# ===== SLIDE 2 - OVERVIEW =====
s = add_slide()
add_title(s, "Ringkasan Area yang Membutuhkan Konfirmasi MKP")
items = [
    "1. dim_date — konvensi kalender (first day of week & tanggal awal generate)",
    "2. Filter business logic SP — scope data backlog, tipe MO & PM activity type",
    "3. Mapping SAP user status — kode status & definisi step",
    "4. config_image.csv — gambar/foto referensi section/unit",
    "5. aging_category.csv — threshold aging & priority",
]
box = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
tf = box.text_frame
tf.word_wrap = True
for i, it in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(14)
    r = p.add_run()
    r.text = it
    r.font.size = Pt(20)
    r.font.color.rgb = DARK
add_footer(s, "1 / 7")


# ===== SLIDE 3 - dim_date =====
s = add_slide()
add_title(s, "1. dim_date — Konvensi Kalender", "Pipeline generate dim_date parquet baru untuk MKP")

box = s.shapes.add_textbox(Inches(0.85), Inches(1.3), Inches(11.9), Inches(0.6))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Fungsi dim_date: tabel dimensi tanggal yang dipakai SP & PBI untuk join dan agregasi berdasarkan waktu (date_id, week_id, year, month_id) — dasar untuk menghitung tren backlog per minggu/bulan/tahun di dashboard."
r.font.italic = True
r.font.size = Pt(13)
r.font.color.rgb = GRAY

add_item(s, 2.1,
    "Konvensi hari pertama minggu (parameter @first_day_of_week) — mis. Senin, Kamis, atau Sabtu seperti BUMA ID.",
    "Kolom week_id (format YYYYWW) dihitung berdasarkan hari pertama minggu ini. Kalau konvensinya beda dengan kalender internal MKP, nomor minggu yang tampil di dashboard tidak akan konsisten dengan kalender kerja mereka.")
add_item(s, 3.6,
    "Tanggal awal data dim_date mulai di-generate (parameter @first_date_of_date_generation). Di BUMA ID sebelumnya hardcode '2020-01-01'.",
    "Menentukan rentang tanggal yang tersedia untuk join. Kalau rentangnya tidak mencakup data backlog historis MKP, baris-baris data tersebut tidak akan punya tanggal pasangan dan bisa hilang dari laporan.")
add_footer(s, "2 / 7")


# ===== SLIDE 4 - Filter business logic =====
s = add_slide()
add_title(s, "2. Filter Business Logic SP — Scope Data Backlog", "Parameterisasi filter yang sebelumnya hardcode untuk BUMA ID")
add_item(s, 1.55,
    "Apakah filter created_by (@moCreatedBy) dan tipe MO (@moType) di MKP hanya bernilai satu (single value, seperti 'DATACOM' & 'MT01' di BUMA ID), atau ada beberapa nilai sekaligus.",
    "Di SP saat ini, kedua parameter ini didesain untuk single value. Jika MKP punya lebih dari satu nilai untuk masing-masing filter, SP perlu didesain ulang agar bisa menerima multi-value (mis. comma-separated atau table-valued parameter) — desain teknisnya belum diputuskan.")
add_item(s, 3.55,
    "Apakah seluruh data backlog MKP boleh ditampilkan tanpa filter created_by, atau hanya subset tertentu yang relevan.",
    "Berkaitan dengan volume data (performa query/laporan jika semua data ditarik) dan sensitivitas data (apakah semua MO layak ditampilkan di dashboard backlog atau ada yang perlu disaring).")
add_item(s, 5.55,
    "Nilai @replacementStatusLvl1 (PM Activity Type di SAP) untuk MKP. Di BUMA ID memakai 'BEX'; di MKP kemungkinan 'BKG', atau bisa jadi ada activity type lain yang juga relevan (multi-value).",
    "Parameter ini memfilter MO berdasarkan PM Activity Type sebelum dihitung sebagai backlog. Kalau activity type MKP berbeda atau lebih dari satu, filter ini perlu disesuaikan (termasuk kemungkinan multi-value) agar data backlog MKP tidak salah ter-filter.")
add_footer(s, "3 / 7")


# ===== SLIDE 5 - SAP user status =====
s = add_slide()
add_title(s, "3. Mapping SAP User Status", "Config mapping baru: config_mapping_sap_user_status.csv")
add_item(s, 1.55,
    "Nilai aktual SAP user status code di data MKP, dan definisi step status MO MKP — mana yang dianggap Closed, In Progress, Waiting Approval, dst (dikonfirmasi satu per satu).",
    "SP butuh mapping dari kode status SAP MKP ke kategori standar (status_category) untuk menghitung backlog per status. Kode status MKP pasti berbeda dari BUMA ID, dan standar MO SAP tidak punya workflow approval default — Waiting Approval adalah konfigurasi tambahan khusus BUMA ID — sehingga step status MKP perlu dikonfirmasi satu per satu.")

box = s.shapes.add_textbox(Inches(0.85), Inches(3.0), Inches(11.9), Inches(0.4))
tf = box.text_frame
r = tf.paragraphs[0].add_run()
r.text = "Config yang dibutuhkan — config_mapping_sap_user_status.csv:"
r.font.bold = True
r.font.size = Pt(14)
r.font.color.rgb = DARK

add_table(s, 0.85, 3.5, 11.9, 3.2,
    ["sap_status (kode SAP — diisi dari data MKP)", "status_category (standar Digiman+)"],
    [
        ["?", "closed"],
        ["?", "approved"],
        ["?", "waiting_approval"],
        ["?", "waiting_planning"],
        ["?", "waiting_part"],
        ["?", "waiting_execution"],
    ])
add_footer(s, "4 / 7")


# ===== SLIDE 6 - config_image.csv =====
s = add_slide()
add_title(s, "4. config_image.csv — Gambar Referensi Section/Unit", "Config baru, disimpan sebagai base64 (chunked)")
add_item(s, 1.6,
    "Untuk site MKP: ada berapa section/unit yang perlu ditampilkan di dashboard, dan untuk masing-masing section tersebut gambar/foto referensi seperti apa yang diinginkan (foto unit aktual, diagram, ikon, dll).",
    "Daftar nama section sendiri sudah mengikuti standar internal Digiman+ (tidak perlu dikonfirmasi ke MKP). Namun config_image.csv menyimpan gambar referensi per section sebagai base64 — supaya config ini bisa dibuat untuk MKP, perlu tahu jumlah section/unit di site MKP dan bentuk gambar yang ingin ditampilkan untuk masing-masing.", accent=True)
add_footer(s, "5 / 7")


# ===== SLIDE 7 - aging_category.csv =====
s = add_slide()
add_title(s, "5. aging_category.csv — Threshold Aging & Priority", "Struktur CSV diperluas: tambah kolom min_days, max_days, priority")
add_item(s, 1.55,
    "Threshold hari (min_days, max_days) per aging_category, dan mapping priority (P1/P2/P3) untuk masing-masing — apakah sama dengan konfigurasi BUMA ID di bawah, atau perlu disesuaikan dengan kebijakan/SLA MKP.",
    "aging_category menunjukkan seberapa lama backlog sudah open (menentukan urgensinya), dan priority dipakai untuk highlight/sorting di dashboard. Keduanya bersifat independen (priority tidak otomatis proporsional dari aging range) sehingga perlu eksplisit dikonfirmasi, bukan otomatis memakai default BUMA ID.")

box = s.shapes.add_textbox(Inches(0.85), Inches(3.0), Inches(11.9), Inches(0.4))
tf = box.text_frame
r = tf.paragraphs[0].add_run()
r.text = "Konfigurasi saat ini (default BUMA ID) — aging_category.csv:"
r.font.bold = True
r.font.size = Pt(14)
r.font.color.rgb = DARK

add_table(s, 0.85, 3.5, 11.9, 2.4,
    ["id", "aging_category", "min_days", "max_days", "priority"],
    [
        ["1", "1-14", "0", "14", "P1"],
        ["2", "15-45", "15", "45", "P2"],
        ["3", ">45", "46", "null (tidak ada batas atas)", "P3"],
    ])
add_footer(s, "6 / 7")


# ===== SLIDE 8 - CLOSING =====
s = add_slide()
add_title(s, "Next Steps")
box = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Setelah ke-5 poin konfirmasi di atas didapatkan dari tim MKP, tim Digiman+ akan melanjutkan:"
r.font.size = Pt(18)
r.font.color.rgb = DARK
next_items = [
    "Membangun pipeline generate dim_date & dim_equipment parquet untuk datamart MKP",
    "Memparameterisasi SP usp_iams_backlog_monitoring sesuai konfigurasi tenant MKP",
    "Menyiapkan config mapping file (SAP user status, equipment category, image, aging category) untuk MKP",
]
for it in next_items:
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.level = 1
    r = p.add_run()
    r.text = "•  " + it
    r.font.size = Pt(16)
    r.font.color.rgb = GRAY
add_footer(s, "7 / 7")

out_path = r"D:\OneDrive-Btech\OneDrive - Bukit Technology\Notebooks\office-notes\digiman+\report\backlog-monitoring\backlog-monitoring-mkp-konfirmasi.pptx"
prs.save(out_path)
print("Saved:", out_path)
